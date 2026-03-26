"""
Gera apresentação profissional do SUAV Bot
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Cores da marca SUAV
ROSA = RGBColor(0xE9, 0x5D, 0x87)       # rosa principal
ROSA_DARK = RGBColor(0xC4, 0x3D, 0x6B)   # rosa escuro
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_ESCURO = RGBColor(0x2D, 0x2D, 0x2D)
CINZA_MEDIO = RGBColor(0x5A, 0x5A, 0x5A)
CINZA_CLARO = RGBColor(0xF5, 0xF0, 0xF0)
VERDE = RGBColor(0x27, 0xAE, 0x60)
AZUL = RGBColor(0x34, 0x98, 0xDB)
DOURADO = RGBColor(0xD4, 0xA5, 0x37)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    """Fundo sólido"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    """Retângulo colorido"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        sld = shape.fill._fill
        solid = sld.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solid is not None:
            srgb = solid.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                a = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                a.set('val', str(int(alpha * 1000)))
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=CINZA_ESCURO, bold=False, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Caixa de texto"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=CINZA_ESCURO, icon=""):
    """Lista com bullets"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon} {item}" if icon else f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)
    return txBox


def add_card(slide, left, top, width, height, title, items, accent_color=ROSA, icon=""):
    """Card com título e lista"""
    # Fundo do card
    card = add_shape(slide, left, top, width, height, BRANCO)
    card.shadow.inherit = False

    # Barra de acento
    add_shape(slide, left, top, Inches(0.06), height, accent_color)

    # Título
    add_text(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.5), Inches(0.4),
             f"{icon}  {title}" if icon else title, font_size=16, color=accent_color, bold=True)

    # Items
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.55), width - Inches(0.5), height - Inches(0.7),
                    items, font_size=13, color=CINZA_MEDIO, icon="•")


def add_stat_box(slide, left, top, width, number, label, color=ROSA):
    """Caixa de estatística"""
    add_shape(slide, left, top, width, Inches(1.2), BRANCO)
    add_text(slide, left, top + Inches(0.1), width, Inches(0.6),
             number, font_size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.7), width, Inches(0.4),
             label, font_size=12, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1 — CAPA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BRANCO)

# Faixa rosa no topo
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

# Linha decorativa
add_shape(slide, Inches(4.5), Inches(2.2), Inches(4.3), Inches(0.04), ROSA)

add_text(slide, Inches(2), Inches(2.4), Inches(9.3), Inches(1.2),
         "SUAV Beauty", font_size=56, color=ROSA, bold=True, align=PP_ALIGN.CENTER,
         font_name='Segoe UI Light')

add_text(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(0.8),
         "Assistente Virtual Inteligente", font_size=28, color=CINZA_ESCURO, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.6),
         "Atendimento automatizado por WhatsApp e Instagram", font_size=18, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# Linha decorativa inferior
add_shape(slide, Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.04), ROSA)

add_text(slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.5),
         "Desenvolvido por Virgílio  |  2026", font_size=14, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — O QUE É A MARI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Conheça a Mari", font_size=38, color=ROSA, bold=True)
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Sua atendente virtual que fala como gente de verdade", font_size=18, color=CINZA_MEDIO)

# Cards
add_card(slide, Inches(0.8), Inches(1.9), Inches(3.6), Inches(2.3),
         "Quem é a Mari?", [
             "Atendente virtual da SUAV",
             "Fala como pessoa real no WhatsApp",
             "Gentil, simpática e acolhedora",
             "Nunca parece um robô",
         ], ROSA, "👩")

add_card(slide, Inches(4.8), Inches(1.9), Inches(3.6), Inches(2.3),
         "Como funciona?", [
             "Responde 24/7 no WhatsApp",
             "Atende pelo Instagram DM",
             "Identifica clientes, profissionais e admin",
             "Agenda, cancela e reagenda sozinha",
         ], AZUL, "⚙️")

add_card(slide, Inches(8.8), Inches(1.9), Inches(3.6), Inches(2.3),
         "Diferenciais", [
             "Sem frases robóticas",
             "Quebra mensagens como gente real",
             "Emojis com moderação natural",
             "Encaminha reclamações pra gerente",
         ], VERDE, "✨")

# Stat boxes
add_stat_box(slide, Inches(1.3), Inches(4.8), Inches(2.5), "101", "Serviços cadastrados", ROSA)
add_stat_box(slide, Inches(4.3), Inches(4.8), Inches(2.5), "24/7", "Atendimento contínuo", AZUL)
add_stat_box(slide, Inches(7.3), Inches(4.8), Inches(2.5), "3", "Perfis de acesso", VERDE)
add_stat_box(slide, Inches(10.3), Inches(4.8), Inches(2.5), "5", "Categorias", DOURADO)


# ============================================================
# SLIDE 3 — FUNCIONALIDADES CLIENTE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Para Clientes", font_size=38, color=ROSA, bold=True)
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Tudo que a cliente consegue fazer pelo WhatsApp e Instagram", font_size=18, color=CINZA_MEDIO)

add_card(slide, Inches(0.5), Inches(1.9), Inches(3.9), Inches(2.5),
         "Informações", [
             "Consultar serviços e preços",
             "Ver profissionais por serviço",
             "Horário de funcionamento",
             "Endereço e estacionamento",
             "Formas de pagamento",
         ], ROSA, "📋")

add_card(slide, Inches(4.7), Inches(1.9), Inches(3.9), Inches(2.5),
         "Agendamento", [
             "Agendar com data, hora e profissional",
             "Ver agendamentos futuros",
             "Cancelar agendamento",
             "Reagendar para outro horário",
             "Verificar disponibilidade em tempo real",
         ], AZUL, "📅")

add_card(slide, Inches(8.9), Inches(1.9), Inches(3.9), Inches(2.5),
         "Experiência", [
             "Lembra o nome da cliente",
             "Salva profissional preferida",
             "Não pede telefone no WhatsApp",
             "Responde em português sempre",
             "Reclamação vai direto pra gerente",
         ], VERDE, "💬")

# Exemplo de conversa
add_shape(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3), CINZA_CLARO)
add_text(slide, Inches(0.8), Inches(4.85), Inches(5), Inches(0.3),
         "Exemplo de conversa real:", font_size=14, color=ROSA, bold=True)

convo = [
    ('Cliente:', '"oi, quanto custa unha gel?"', CINZA_ESCURO),
    ('Mari:', '"Manutenção unha gel a partir de R$149.', ROSA_DARK),
    ('', 'Quer que eu veja um horário pra você?"', ROSA_DARK),
    ('Cliente:', '"sim, sexta à tarde com a Larissa"', CINZA_ESCURO),
    ('Mari:', '"A Larissa tem vaga às 14h e 16h na sexta.', ROSA_DARK),
    ('', 'Qual vc prefere?"', ROSA_DARK),
]
y = 5.25
for who, text, color in convo:
    full = f"{who} {text}" if who else f"            {text}"
    add_text(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.28),
             full, font_size=13, color=color, bold=(who != ''))
    y += 0.27


# ============================================================
# SLIDE 4 — FUNCIONALIDADES ADMIN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Para a Gerente / Dona", font_size=38, color=ROSA, bold=True)
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "24 funções de gestão completa pelo WhatsApp — sem precisar abrir sistema", font_size=18, color=CINZA_MEDIO)

add_card(slide, Inches(0.5), Inches(1.9), Inches(2.9), Inches(2.8),
         "Consultas", [
             "Faturamento por período",
             "Agenda do dia",
             "Estatísticas de agendamento",
             "Ranking de profissionais",
             "Histórico de clientes",
             "Buscar cliente",
         ], ROSA, "📊")

add_card(slide, Inches(3.7), Inches(1.9), Inches(2.9), Inches(2.8),
         "Analytics", [
             "Relatório de no-shows",
             "Análise de cancelamentos",
             "Horários de pico",
             "Clientes inativas (retenção)",
             "Relatório diário automático",
         ], AZUL, "📈")

add_card(slide, Inches(6.9), Inches(1.9), Inches(2.9), Inches(2.8),
         "Gestão", [
             "Agendar em nome de cliente",
             "Ativar/desativar profissional",
             "Alterar preços",
             "Alterar horários de trabalho",
             "Bloquear/desbloquear horários",
         ], VERDE, "🔧")

add_card(slide, Inches(10.1), Inches(1.9), Inches(2.9), Inches(2.8),
         "Reclamações", [
             "Recebe alerta de reclamação",
             "Bot para de responder cliente",
             "Listar reclamações pendentes",
             "Resolver escalação",
             "Bot volta a atender",
         ], DOURADO, "⚠️")

# Destaque
add_shape(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.6), CINZA_CLARO)
add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
         "Exemplos de comandos naturais:", font_size=14, color=ROSA, bold=True)

cmds = [
    '"faturamento desse mês"  →  Mostra total + por profissional + por categoria',
    '"ranking das profissionais"  →  Top performers por receita ou quantidade',
    '"clientes sumidas há 60 dias"  →  Lista quem não volta, com último serviço e data',
    '"faltas desse mês"  →  No-shows, prejuízo total, clientes reincidentes',
]
add_bullet_list(slide, Inches(1), Inches(5.6), Inches(11), Inches(1.1),
                cmds, font_size=13, color=CINZA_MEDIO, icon="💬")


# ============================================================
# SLIDE 5 — FUNCIONALIDADES PROFISSIONAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Para Profissionais", font_size=38, color=ROSA, bold=True)
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Cada profissional conversa com a Mari pelo seu próprio WhatsApp", font_size=18, color=CINZA_MEDIO)

add_card(slide, Inches(0.5), Inches(1.9), Inches(3.9), Inches(3.0),
         "Agenda Pessoal", [
             "Ver agenda do dia",
             "Ver agenda da semana",
             "Próxima cliente (olhada rápida)",
             "Marcar atendimento como concluído",
             "Bloquear horário (almoço, médico)",
             "Desbloquear horário",
         ], ROSA, "📅")

add_card(slide, Inches(4.7), Inches(1.9), Inches(3.9), Inches(3.0),
         "Faturamento Pessoal", [
             "Faturamento de hoje",
             "Faturamento da semana",
             "Faturamento do mês",
             "Ticket médio",
             "Número de atendimentos",
         ], AZUL, "💰")

add_card(slide, Inches(8.9), Inches(1.9), Inches(3.9), Inches(3.0),
         "Notificações Automáticas", [
             "Novo agendamento (serviço + cliente + horário + valor)",
             "Cancelamento de cliente",
             "Remarcação de horário",
             "Lembrete 24h antes",
             "Lembrete 1h antes",
         ], VERDE, "🔔")

# Segurança
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5), CINZA_CLARO)
add_text(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.4),
         "🔒  Segurança e Privacidade", font_size=16, color=ROSA, bold=True)

security = [
    "Profissional só vê SUA agenda — nunca dados de outras",
    "Faturamento do salão → direcionado para a gerente",
    "Dados de clientes → somente informações do atendimento dela",
    "Cada profissional é identificada automaticamente pelo número do WhatsApp",
]
add_bullet_list(slide, Inches(1), Inches(5.8), Inches(11), Inches(1.0),
                security, font_size=13, color=CINZA_MEDIO, icon="✅")


# ============================================================
# SLIDE 6 — NOTIFICAÇÕES E LEMBRETES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Notificações Inteligentes", font_size=38, color=ROSA, bold=True)
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Tudo automático — ninguém precisa lembrar de nada", font_size=18, color=CINZA_MEDIO)

# Timeline
steps = [
    ("📅", "Cliente agenda", "Mari confirma e avisa\na profissional + gerente", ROSA),
    ("⏰", "24h antes", "Lembrete automático\npara a cliente", AZUL),
    ("🔔", "1h antes", "Segundo lembrete\ncom endereço", VERDE),
    ("✅", "Após atendimento", "Profissional marca\ncomo concluído", DOURADO),
    ("📊", "Fim do dia", "Relatório diário\npara a gerente", ROSA_DARK),
]

for i, (icon, title, desc, color) in enumerate(steps):
    x = Inches(0.8 + i * 2.5)
    y = Inches(2.0)

    # Círculo
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    add_text(slide, x + Inches(0.6), y + Inches(0.15), Inches(0.9), Inches(0.6),
             icon, font_size=28, color=BRANCO, align=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(1.1), Inches(2.1), Inches(0.4),
             title, font_size=15, color=color, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(1.5), Inches(2.1), Inches(0.7),
             desc, font_size=12, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# Linha conectando
add_shape(slide, Inches(1.65), Inches(2.43), Inches(10.05), Inches(0.03), CINZA_CLARO)

# Exemplo de notificação
add_shape(slide, Inches(1.5), Inches(4.5), Inches(5), Inches(2.3), CINZA_CLARO)
add_text(slide, Inches(1.8), Inches(4.6), Inches(4.5), Inches(0.3),
         "Notificação para profissional:", font_size=13, color=ROSA, bold=True)

notif_lines = [
    "Nova cliente agendada! ✨",
    "",
    "Serviço: Manutenção Unha Gel",
    "Cliente: Camila (27) 99999-0000",
    "Data: sexta-feira, 28/03/2026",
    "Horário: 14:00 às 15:30",
    "Valor: R$ 149,00",
]
y = 4.95
for line in notif_lines:
    add_text(slide, Inches(1.8), Inches(y), Inches(4.5), Inches(0.22),
             line, font_size=12, color=CINZA_ESCURO if line else CINZA_MEDIO)
    y += 0.2

# Lembrete para cliente
add_shape(slide, Inches(7), Inches(4.5), Inches(5), Inches(2.3), CINZA_CLARO)
add_text(slide, Inches(7.3), Inches(4.6), Inches(4.5), Inches(0.3),
         "Lembrete para cliente (24h antes):", font_size=13, color=AZUL, bold=True)

reminder_lines = [
    "Oi Camila! 😊",
    "",
    "Só passando pra lembrar do seu",
    "agendamento amanhã:",
    "",
    "💅 Manutenção Unha Gel com a Larissa",
    "📍 R. Goiânia, 234 - Itapoã, Vila Velha",
    "⏰ Sexta, 28/03 às 14:00",
]
y = 4.95
for line in reminder_lines:
    add_text(slide, Inches(7.3), Inches(y), Inches(4.5), Inches(0.22),
             line, font_size=12, color=CINZA_ESCURO if line else CINZA_MEDIO)
    y += 0.2


# ============================================================
# SLIDE 7 — FLUXO DE RECLAMAÇÃO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Gestão de Reclamações", font_size=38, color=ROSA, bold=True)
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Bot nunca tenta resolver sozinho — sempre escala para a gerente", font_size=18, color=CINZA_MEDIO)

flow_steps = [
    ("1", "Cliente reclama", "\"Minha unha quebrou\nno mesmo dia\"", ROSA),
    ("2", "Mari encaminha", "Chama report_complaint\ne avisa a cliente", AZUL),
    ("3", "Gerente recebe", "WhatsApp com detalhes\nda reclamação", VERDE),
    ("4", "Bot silencia", "Para de responder\naquela cliente", DOURADO),
    ("5", "Gerente resolve", "Diz \"resolve escalação\"\ne bot volta", ROSA_DARK),
]

for i, (num, title, desc, color) in enumerate(flow_steps):
    x = Inches(0.6 + i * 2.5)
    y = Inches(2.2)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    add_text(slide, x + Inches(0.65), y + Inches(0.12), Inches(0.8), Inches(0.6),
             num, font_size=30, color=BRANCO, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(1.0), Inches(2.1), Inches(0.4),
             title, font_size=15, color=color, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(1.4), Inches(2.1), Inches(0.7),
             desc, font_size=12, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(1.45), Inches(2.58), Inches(10.4), Inches(0.03), CINZA_CLARO)

# Destaque de segurança
add_shape(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(2), CINZA_CLARO)
add_text(slide, Inches(2), Inches(4.95), Inches(9), Inches(0.4),
         "Por que isso é importante?", font_size=16, color=ROSA, bold=True)

reasons = [
    "Evita que o bot dê respostas inadequadas em situações delicadas",
    "Gerente tem controle total sobre a resolução",
    "Cliente sabe que uma pessoa real vai cuidar do caso",
    "Todas as mensagens da cliente ficam salvas no histórico para contexto",
    "Bot volta a funcionar normalmente após resolução",
]
add_bullet_list(slide, Inches(2), Inches(5.35), Inches(9), Inches(1.3),
                reasons, font_size=13, color=CINZA_MEDIO, icon="✅")


# ============================================================
# SLIDE 8 — TECNOLOGIA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Tecnologia", font_size=38, color=ROSA, bold=True)
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Arquitetura robusta e escalável", font_size=18, color=CINZA_MEDIO)

add_card(slide, Inches(0.5), Inches(1.9), Inches(3.9), Inches(2.5),
         "Inteligência Artificial", [
             "OpenAI GPT-4.1 Nano",
             "Function calling com JSON Schema",
             "3 system prompts especializados",
             "Anti-alucinação de nomes",
             "Filtro anti-bot automático",
         ], ROSA, "🧠")

add_card(slide, Inches(4.7), Inches(1.9), Inches(3.9), Inches(2.5),
         "Infraestrutura", [
             "TypeScript / Node.js",
             "Prisma ORM + banco relacional",
             "WhatsApp via Baileys",
             "Instagram Graph API",
             "Cron jobs para lembretes",
         ], AZUL, "🏗️")

add_card(slide, Inches(8.9), Inches(1.9), Inches(3.9), Inches(2.5),
         "Experiência do Usuário", [
             "Typing indicators realistas",
             "Mensagens quebradas como gente real",
             "Rate limiting anti-spam",
             "Detecção de mídia",
             "Histórico persistente",
         ], VERDE, "💎")

# Detalhes técnicos
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.1), CINZA_CLARO)
add_text(slide, Inches(1), Inches(4.9), Inches(11), Inches(0.4),
         "Detalhes técnicos", font_size=14, color=ROSA, bold=True)

tech_details = [
    "34 cenários de teste automatizados com 46 checks — 100% de aprovação",
    "3 perfis de acesso (cliente, profissional, admin) com funções isoladas por papel",
    "Identificação automática por número de telefone — sem login necessário",
    "Suporte a WhatsApp e Instagram com adaptação automática de comportamento",
    "Notificações em tempo real para profissionais, clientes e administradores",
    "Sistema de lembretes com prevenção de duplicatas via flags no banco",
]
add_bullet_list(slide, Inches(1), Inches(5.3), Inches(11), Inches(1.5),
                tech_details, font_size=13, color=CINZA_MEDIO, icon="→")


# ============================================================
# SLIDE 9 — CONTATO / ENCERRAMENTO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ROSA)

add_shape(slide, Inches(0), Inches(2.5), W, Inches(3), CINZA_CLARO)

add_text(slide, Inches(2), Inches(2.8), Inches(9.3), Inches(1),
         "SUAV Beauty Bot", font_size=48, color=ROSA, bold=True, align=PP_ALIGN.CENTER,
         font_name='Segoe UI Light')

add_text(slide, Inches(2), Inches(3.7), Inches(9.3), Inches(0.6),
         "Transformando atendimento em experiência", font_size=22, color=CINZA_ESCURO, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.5),
         "101 serviços  •  24/7  •  WhatsApp + Instagram  •  3 perfis de acesso", font_size=16, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.5), Inches(5.5), Inches(4.3), Inches(0.04), ROSA)

add_text(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(0.5),
         "Desenvolvido por Virgílio", font_size=18, color=CINZA_ESCURO, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(6.3), Inches(9.3), Inches(0.4),
         "📞 (96) 99204-8681   |   📍 Vila Velha, ES", font_size=15, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)


# ============================================================
# SALVAR
# ============================================================
output_path = r'C:\Users\Virgilio\Desktop\SUAV_Bot_Apresentacao.pptx'
prs.save(output_path)
print(f'Apresentação salva em: {output_path}')
print(f'Total de slides: {len(prs.slides)}')
